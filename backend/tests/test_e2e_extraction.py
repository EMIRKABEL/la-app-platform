"""End-to-end test: course → unit → lesson → curriculum upload → extraction.

This test exercises the full extraction workflow:
1. Create a Course
2. Create a Unit inside the Course
3. Create a Lesson inside the Unit
4. Upload a real PPTX curriculum file (generated with python-pptx)
5. Extract the PPTX structure via POST /api/curriculum/{id}/extract
6. Verify structured JSON is stored
7. Verify GET /api/curriculum/{id}/extraction returns it
"""

import io
import os
import tempfile
import shutil

from pptx import Presentation


def test_e2e_curriculum_extraction(client, monkeypatch):
    """Full end-to-end curriculum extraction workflow."""
    tmp_storage = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp_storage)

    from app.core import config
    config.get_settings.cache_clear()

    try:
        # 1. Create Course
        r = client.post(
            "/api/courses",
            json={"name": "E2E Extract Course", "description": "Extraction e2e"},
        )
        assert r.status_code == 201
        course_id = r.json()["id"]

        # 2. Create Unit
        r2 = client.post(
            f"/api/courses/{course_id}/units",
            json={"number": 1, "title": "Unit 1"},
        )
        assert r2.status_code == 201
        unit_id = r2.json()["id"]

        # 3. Create Lesson
        r3 = client.post(
            f"/api/units/{unit_id}/lessons",
            json={"number": 1, "title": "Lesson 1"},
        )
        assert r3.status_code == 201
        lesson_id = r3.json()["id"]

        # 4. Generate and upload a real PPTX file
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Lesson 1"
        slide.placeholders[1].text = "Hello World"

        # Add a second slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Slide 2"
        slide2.placeholders[1].text = "Second slide content"

        pptx_buf = io.BytesIO()
        prs.save(pptx_buf)
        pptx_content = pptx_buf.getvalue()

        r4 = client.post(
            f"/api/lessons/{lesson_id}/curriculum",
            files={
                "file": (
                    "e2e-extract.pptx",
                    io.BytesIO(pptx_content),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
        assert r4.status_code == 201
        source_id = r4.json()["id"]
        assert r4.json()["processing_status"] == "pending"

        # 5. Extract the PPTX
        r5 = client.post(f"/api/curriculum/{source_id}/extract")
        assert r5.status_code == 200
        extract_body = r5.json()

        assert extract_body["processing_status"] == "completed"
        assert extract_body["extracted_data"] is not None
        assert extract_body["extracted_data"]["source_type"] == "pptx"
        assert extract_body["extracted_data"]["metadata"]["slide_count"] == 2

        slides = extract_body["extracted_data"]["slides"]
        assert len(slides) == 2
        assert slides[0]["slide_number"] == 1
        assert slides[0]["title"] == "Lesson 1"
        assert "Hello World" in slides[0]["texts"]
        assert slides[1]["slide_number"] == 2
        assert slides[1]["title"] == "Slide 2"

        # 6. Verify GET extraction returns the saved data
        r6 = client.get(f"/api/curriculum/{source_id}/extraction")
        assert r6.status_code == 200
        get_body = r6.json()

        assert get_body["id"] == source_id
        assert get_body["processing_status"] == "completed"
        assert get_body["extracted_data"] is not None
        assert get_body["extracted_data"]["source_type"] == "pptx"
        assert len(get_body["extracted_data"]["slides"]) == 2

    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp_storage, ignore_errors=True)
