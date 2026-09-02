"""Tests for course-level curriculum upload, listing, and extraction.

Test PPTX files are generated in-memory using ``python-pptx``.
"""

import io
import os
import tempfile
import shutil
import uuid

from pptx import Presentation


# ── Helpers ──────────────────────────────────────────────────────


def _create_course(client, name="Test Course"):
    """Create a course and return its id."""
    resp = client.post(
        "/api/courses",
        json={"name": name, "description": "For course curriculum tests"},
    )
    assert resp.status_code == 201
    return resp.json()["id"]


def _make_pptx(title="Slide Title", text="Hello"):
    """Generate a minimal PPTX in memory and return its bytes."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_file_tuple(filename, content, mime="application/octet-stream"):
    """Return a files dict suitable for the client.post files= parameter."""
    return {"files": (filename, io.BytesIO(content), mime)}


def _pptx_mime():
    return "application/vnd.openxmlformats-officedocument.presentationml.presentation"


# ── Tests: Multi-file upload ─────────────────────────────────────


def test_upload_multiple_files(client, monkeypatch):
    """POST /api/courses/{id}/curriculum with multiple files should create all."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)

        pptx1 = _make_pptx("Slide 1", "Content 1")
        pptx2 = _make_pptx("Slide 2", "Content 2")

        resp = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[
                ("files", ("first.pptx", io.BytesIO(pptx1), _pptx_mime())),
                ("files", ("second.pptx", io.BytesIO(pptx2), _pptx_mime())),
            ],
        )
        assert resp.status_code == 201
        body = resp.json()
        assert len(body) == 2
        filenames = [r["original_filename"] for r in body]
        assert "first.pptx" in filenames
        assert "second.pptx" in filenames
        for r in body:
            assert r["course_id"] == course_id
            assert r["file_type"] == "pptx"
            assert r["processing_status"] == "pending"
            assert r["storage_path"].startswith("curriculum/courses/")
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_upload_single_file(client, monkeypatch):
    """POST with a single file should also work."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        pptx_bytes = _make_pptx("Single", "Single content")

        resp = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[("files", ("single.pptx", io.BytesIO(pptx_bytes), _pptx_mime()))],
        )
        assert resp.status_code == 201
        body = resp.json()
        assert len(body) == 1
        assert body[0]["original_filename"] == "single.pptx"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


# ── Tests: Listing ───────────────────────────────────────────────


def test_list_course_curriculum(client, monkeypatch):
    """GET /api/courses/{id}/curriculum returns all uploaded files."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        pptx1 = _make_pptx("A", "a")
        pptx2 = _make_pptx("B", "b")

        client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[
                ("files", ("a.pptx", io.BytesIO(pptx1), _pptx_mime())),
                ("files", ("b.pdf", io.BytesIO(b"%PDF-1.4 fake"), "application/pdf")),
            ],
        )

        resp = client.get(f"/api/courses/{course_id}/curriculum")
        assert resp.status_code == 200
        body = resp.json()
        assert len(body) == 2
        filenames = [r["original_filename"] for r in body]
        assert "a.pptx" in filenames
        assert "b.pdf" in filenames
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_list_course_curriculum_empty(client):
    """GET curriculum on a course with no uploads returns []."""
    course_id = _create_course(client)
    resp = client.get(f"/api/courses/{course_id}/curriculum")
    assert resp.status_code == 200
    assert resp.json() == []


# ── Tests: Extraction ────────────────────────────────────────────


def test_extract_pptx(client, monkeypatch):
    """POST /api/course-curriculum/{id}/extract extracts PPTX."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        pptx_bytes = _make_pptx("Lesson 1", "Hello World")

        upload_resp = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[("files", ("extract.pptx", io.BytesIO(pptx_bytes), _pptx_mime()))],
        )
        assert upload_resp.status_code == 201
        source_id = upload_resp.json()[0]["id"]

        # Extract
        resp = client.post(f"/api/course-curriculum/{source_id}/extract")
        assert resp.status_code == 200
        body = resp.json()
        assert body["processing_status"] == "completed"
        assert body["extracted_data"]["source_type"] == "pptx"
        assert body["extracted_data"]["metadata"]["slide_count"] == 1
        assert body["extracted_data"]["slides"][0]["title"] == "Lesson 1"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_get_extraction(client, monkeypatch):
    """GET /api/course-curriculum/{id}/extraction returns saved extraction."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        pptx_bytes = _make_pptx("Titled", "Body text")

        upload_resp = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[("files", ("get.pptx", io.BytesIO(pptx_bytes), _pptx_mime()))],
        )
        source_id = upload_resp.json()[0]["id"]

        # Extract first
        client.post(f"/api/course-curriculum/{source_id}/extract")

        # Now GET
        resp = client.get(f"/api/course-curriculum/{source_id}/extraction")
        assert resp.status_code == 200
        body = resp.json()
        assert body["id"] == source_id
        assert body["processing_status"] == "completed"
        assert body["extracted_data"] is not None
        assert body["extracted_data"]["source_type"] == "pptx"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_extract_unsupported_file_type(client, monkeypatch):
    """Extracting a non-PPTX file returns 415."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        upload_resp = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[("files", ("doc.pdf", io.BytesIO(b"%PDF-1.4 fake"), "application/pdf"))],
        )
        assert upload_resp.status_code == 201
        source_id = upload_resp.json()[0]["id"]

        resp = client.post(f"/api/course-curriculum/{source_id}/extract")
        assert resp.status_code == 415
        detail = resp.json()["detail"]
        assert ".pdf" in detail
        assert ".pptx" in detail
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_extract_not_found(client):
    """POST extract for a non-existent source returns 404."""
    random_id = str(uuid.uuid4())
    resp = client.post(f"/api/course-curriculum/{random_id}/extract")
    assert resp.status_code == 404
    assert resp.json()["detail"] == "Course curriculum source not found"


def test_get_extraction_not_found(client):
    """GET extraction for a non-existent source returns 404."""
    random_id = str(uuid.uuid4())
    resp = client.get(f"/api/course-curriculum/{random_id}/extraction")
    assert resp.status_code == 404
    assert resp.json()["detail"] == "Course curriculum source not found"


# ── Tests: Nonexistent course ────────────────────────────────────


def test_upload_nonexistent_course(client):
    """POST curriculum to a non-existent course returns 404."""
    random_id = str(uuid.uuid4())
    pptx_bytes = _make_pptx()
    resp = client.post(
        f"/api/courses/{random_id}/curriculum",
        files=[("files", ("test.pptx", io.BytesIO(pptx_bytes), _pptx_mime()))],
    )
    assert resp.status_code == 404
    assert resp.json()["detail"] == "Course not found"


def test_list_nonexistent_course(client):
    """GET curriculum for a non-existent course returns []."""
    random_id = str(uuid.uuid4())
    resp = client.get(f"/api/courses/{random_id}/curriculum")
    # No FK enforcement on GET — returns empty list
    assert resp.status_code == 200
    assert resp.json() == []


# ── Tests: Storage paths and filename safety ─────────────────────


def test_storage_path_format(client, monkeypatch):
    """Verify storage path is curriculum/courses/{course_id}/{filename}."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        pptx_bytes = _make_pptx()

        resp = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[("files", ("test.pptx", io.BytesIO(pptx_bytes), _pptx_mime()))],
        )
        assert resp.status_code == 201
        storage_path = resp.json()[0]["storage_path"]
        assert storage_path.startswith(f"curriculum/courses/{course_id}/")
        assert storage_path.endswith("test.pptx")

        # Verify file exists on disk
        full_path = os.path.join(tmp, storage_path.replace("/", os.sep))
        assert os.path.isfile(full_path)
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_filename_safety(client, monkeypatch):
    """Path traversal attempts in filenames must be sanitized."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        pptx_bytes = _make_pptx()

        resp = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[("files", ("../../../etc/passwd.pptx", io.BytesIO(pptx_bytes), _pptx_mime()))],
        )
        assert resp.status_code == 201
        storage_path = resp.json()[0]["storage_path"]
        assert ".." not in storage_path
        assert "/" not in storage_path.split(f"curriculum/courses/{course_id}/")[1]
        # The filename should be just passwd.pptx
        assert storage_path.endswith("passwd.pptx")
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_duplicate_filename_not_overwritten(client, monkeypatch):
    """Uploading the same filename twice should not overwrite."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        pptx1 = _make_pptx("First", "first content")
        pptx2 = _make_pptx("Second", "second content")

        r1 = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[("files", ("dup.pptx", io.BytesIO(pptx1), _pptx_mime()))],
        )
        r2 = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[("files", ("dup.pptx", io.BytesIO(pptx2), _pptx_mime()))],
        )

        path1 = r1.json()[0]["storage_path"]
        path2 = r2.json()[0]["storage_path"]
        assert path1 != path2  # Different storage paths

        # Both files exist
        full1 = os.path.join(tmp, path1.replace("/", os.sep))
        full2 = os.path.join(tmp, path2.replace("/", os.sep))
        assert os.path.isfile(full1)
        assert os.path.isfile(full2)
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


# ── Tests: Unsupported file rejection ────────────────────────────


def test_unsupported_file_skipped(client, monkeypatch):
    """Uploading a .txt file alongside a .pptx should skip the .txt."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        pptx_bytes = _make_pptx()

        resp = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[
                ("files", ("good.pptx", io.BytesIO(pptx_bytes), _pptx_mime())),
                ("files", ("bad.txt", io.BytesIO(b"plain text"), "text/plain")),
            ],
        )
        assert resp.status_code == 201
        body = resp.json()
        # Only the valid file should be created
        assert len(body) == 1
        assert body[0]["original_filename"] == "good.pptx"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


# ── Tests: Processing status ─────────────────────────────────────


def test_status_transitions(client, monkeypatch):
    """Verify processing_status: pending → processing → completed."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        course_id = _create_course(client)
        pptx_bytes = _make_pptx("Status", "status test")

        upload_resp = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[("files", ("status.pptx", io.BytesIO(pptx_bytes), _pptx_mime()))],
        )
        source_id = upload_resp.json()[0]["id"]
        assert upload_resp.json()[0]["processing_status"] == "pending"

        # Extract
        extract_resp = client.post(f"/api/course-curriculum/{source_id}/extract")
        assert extract_resp.status_code == 200
        assert extract_resp.json()["processing_status"] == "completed"

        # GET should also show completed
        get_resp = client.get(f"/api/course-curriculum/{source_id}/extraction")
        assert get_resp.status_code == 200
        assert get_resp.json()["processing_status"] == "completed"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)
