"""E2E test: create course → upload two PPTX → extract one → retrieve extraction.

This test does NOT create any units or lessons — course-level curriculum
is independent of the unit/lesson hierarchy.
"""

import io
import os
import tempfile
import shutil

from pptx import Presentation


def test_e2e_course_curriculum(client, monkeypatch):
    """Full end-to-end course-level curriculum workflow."""
    tmp_storage = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp_storage)

    from app.core import config
    config.get_settings.cache_clear()

    pptx_mime = (
        "application/vnd.openxmlformats-officedocument"
        ".presentationml.presentation"
    )

    try:
        # 1. Create Course "Sprout 1"
        r = client.post(
            "/api/courses",
            json={"name": "Sprout 1", "description": "E2E course curriculum test"},
        )
        assert r.status_code == 201
        course_id = r.json()["id"]
        assert r.json()["name"] == "Sprout 1"

        # 2. Generate two PPTX files
        prs1 = Presentation()
        slide1 = prs1.slides.add_slide(prs1.slide_layouts[0])
        slide1.shapes.title.text = "Lesson 1"
        slide1.placeholders[1].text = "Greetings"
        buf1 = io.BytesIO()
        prs1.save(buf1)
        pptx1_bytes = buf1.getvalue()

        prs2 = Presentation()
        slide2 = prs2.slides.add_slide(prs2.slide_layouts[0])
        slide2.shapes.title.text = "Lesson 2"
        slide2.placeholders[1].text = "Numbers"
        buf2 = io.BytesIO()
        prs2.save(buf2)
        pptx2_bytes = buf2.getvalue()

        # 3. Upload both PPTX files at course level
        r2 = client.post(
            f"/api/courses/{course_id}/curriculum",
            files=[
                ("files", ("sprout-lesson-1.pptx", io.BytesIO(pptx1_bytes), pptx_mime)),
                ("files", ("sprout-lesson-2.pptx", io.BytesIO(pptx2_bytes), pptx_mime)),
            ],
        )
        assert r2.status_code == 201
        uploaded = r2.json()
        assert len(uploaded) == 2

        # Verify both saved
        filenames = [u["original_filename"] for u in uploaded]
        assert "sprout-lesson-1.pptx" in filenames
        assert "sprout-lesson-2.pptx" in filenames

        for u in uploaded:
            assert u["course_id"] == course_id
            assert u["file_type"] == "pptx"
            assert u["processing_status"] == "pending"
            assert u["storage_path"].startswith(f"curriculum/courses/{course_id}/")

        # 4. Verify files on disk
        for u in uploaded:
            full_path = os.path.join(
                tmp_storage, u["storage_path"].replace("/", os.sep)
            )
            assert os.path.isfile(full_path), f"File not found at {full_path}"

        # 5. List course curriculum
        r3 = client.get(f"/api/courses/{course_id}/curriculum")
        assert r3.status_code == 200
        listed = r3.json()
        assert len(listed) == 2

        # 6. Extract the first one
        source_id = uploaded[0]["id"]
        r4 = client.post(f"/api/course-curriculum/{source_id}/extract")
        assert r4.status_code == 200
        extract_body = r4.json()

        assert extract_body["processing_status"] == "completed"
        assert extract_body["extracted_data"] is not None
        assert extract_body["extracted_data"]["source_type"] == "pptx"
        assert extract_body["extracted_data"]["metadata"]["slide_count"] == 1

        slides = extract_body["extracted_data"]["slides"]
        assert slides[0]["title"] == "Lesson 1"
        assert "Greetings" in slides[0]["texts"]

        # 7. Retrieve extraction via GET
        r5 = client.get(f"/api/course-curriculum/{source_id}/extraction")
        assert r5.status_code == 200
        get_body = r5.json()

        assert get_body["id"] == source_id
        assert get_body["processing_status"] == "completed"
        assert get_body["extracted_data"] is not None
        assert get_body["extracted_data"]["source_type"] == "pptx"
        assert len(get_body["extracted_data"]["slides"]) == 1

        # 8. Verify the second file is still pending (not extracted)
        source2_id = uploaded[1]["id"]
        r6 = client.get(f"/api/course-curriculum/{source2_id}/extraction")
        assert r6.status_code == 200
        assert r6.json()["processing_status"] == "pending"
        assert r6.json()["extracted_data"] is None

    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp_storage, ignore_errors=True)
