"""Tests for PPTX curriculum extraction.

Test PPTX files are generated in-memory using ``python-pptx`` and are
not persisted to disk.
"""

import io
import os
import tempfile
import shutil

from pptx import Presentation
from pptx.util import Inches


# ── Helpers ──────────────────────────────────────────────────────


def _create_course_unit_lesson(client):
    """Helper: create course → unit → lesson and return lesson_id."""
    course_resp = client.post(
        "/api/courses",
        json={"name": "Extraction Course", "description": "For extraction tests"},
    )
    course_id = course_resp.json()["id"]
    unit_resp = client.post(
        f"/api/courses/{course_id}/units",
        json={"number": 1, "title": "Unit 1"},
    )
    unit_id = unit_resp.json()["id"]
    lesson_resp = client.post(
        f"/api/units/{unit_id}/lessons",
        json={"number": 1, "title": "Lesson 1"},
    )
    return lesson_resp.json()["id"]


def _make_simple_pptx() -> bytes:
    """Generate a minimal PPTX with one slide containing a title and text."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Lesson 1"
    slide.placeholders[1].text = "Hello"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_multi_slide_pptx() -> bytes:
    """Generate a PPTX with 3 slides, titles, and text blocks."""
    prs = Presentation()

    # Slide 1 — Title + subtitle
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Lesson 1"
    slide1.placeholders[1].text = "Introduction"

    # Slide 2 — Title + content
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Slide 2 Title"
    body = slide2.placeholders[1]
    body.text = "First paragraph"
    p2 = body.text_frame.add_paragraph()
    p2.text = "Second paragraph"

    # Slide 3 — Title + content
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Slide 3 Title"
    slide3.placeholders[1].text = "Content here"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_notes() -> bytes:
    """Generate a PPTX with speaker notes on a slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Titled Slide"
    slide.placeholders[1].text = "Body text"

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are speaker notes"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_table() -> bytes:
    """Generate a PPTX with a table on a slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    slide.shapes.title.text = "Table Slide"

    # Add a 2x3 table
    table_shape = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(4), Inches(2))
    table = table_shape.table

    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(0, 2).text = "C"
    table.cell(1, 0).text = "D"
    table.cell(1, 1).text = "E"
    table.cell(1, 2).text = "F"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _upload_pptx(client, lesson_id: str, filename: str, content: bytes):
    """Upload a PPTX file and return the response body."""
    resp = client.post(
        f"/api/lessons/{lesson_id}/curriculum",
        files={
            "file": (
                filename,
                io.BytesIO(content),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert resp.status_code == 201
    return resp.json()


# ── Tests: PPTX Extraction ──────────────────────────────────────


def test_extract_simple_pptx(client, monkeypatch):
    """Extract a simple 1-slide PPTX and verify structure."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        pptx_bytes = _make_simple_pptx()
        source = _upload_pptx(client, lesson_id, "simple.pptx", pptx_bytes)
        source_id = source["id"]

        # Extract
        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 200
        body = resp.json()

        assert body["id"] == source_id
        assert body["processing_status"] == "completed"
        assert body["extracted_data"] is not None
        assert body["extracted_data"]["source_type"] == "pptx"
        assert body["extracted_data"]["metadata"]["slide_count"] == 1

        slides = body["extracted_data"]["slides"]
        assert len(slides) == 1
        assert slides[0]["slide_number"] == 1
        assert slides[0]["title"] == "Lesson 1"
        assert "Hello" in slides[0]["texts"]
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_extract_multiple_slides(client, monkeypatch):
    """Extract a 3-slide PPTX and verify slide order and count."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        pptx_bytes = _make_multi_slide_pptx()
        source = _upload_pptx(client, lesson_id, "multi.pptx", pptx_bytes)
        source_id = source["id"]

        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 200
        slides = resp.json()["extracted_data"]["slides"]

        assert len(slides) == 3
        assert slides[0]["slide_number"] == 1
        assert slides[1]["slide_number"] == 2
        assert slides[2]["slide_number"] == 3
        # Verify slide order preserved
        assert slides[0]["title"] == "Lesson 1"
        assert slides[1]["title"] == "Slide 2 Title"
        assert slides[2]["title"] == "Slide 3 Title"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_extract_text_blocks(client, monkeypatch):
    """Verify text extraction captures all visible text blocks."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        pptx_bytes = _make_multi_slide_pptx()
        source = _upload_pptx(client, lesson_id, "texts.pptx", pptx_bytes)
        source_id = source["id"]

        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 200
        slides = resp.json()["extracted_data"]["slides"]

        # Slide 2 has two paragraphs in the body
        slide2 = slides[1]
        assert "First paragraph" in slide2["texts"]
        assert "Second paragraph" in slide2["texts"]
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_extract_title(client, monkeypatch):
    """Verify title extraction returns the slide title."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        pptx_bytes = _make_simple_pptx()
        source = _upload_pptx(client, lesson_id, "title.pptx", pptx_bytes)
        source_id = source["id"]

        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 200
        slides = resp.json()["extracted_data"]["slides"]

        assert slides[0]["title"] == "Lesson 1"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_extract_table(client, monkeypatch):
    """Verify table extraction captures table cell contents."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        pptx_bytes = _make_pptx_with_table()
        source = _upload_pptx(client, lesson_id, "table.pptx", pptx_bytes)
        source_id = source["id"]

        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 200
        slides = resp.json()["extracted_data"]["slides"]

        tables = slides[0]["tables"]
        assert len(tables) == 1
        table = tables[0]
        assert len(table) == 2  # 2 rows
        assert len(table[0]) == 3  # 3 columns
        assert table[0][0] == "A"
        assert table[1][2] == "F"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_extract_speaker_notes(client, monkeypatch):
    """Verify speaker notes extraction."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        pptx_bytes = _make_pptx_with_notes()
        source = _upload_pptx(client, lesson_id, "notes.pptx", pptx_bytes)
        source_id = source["id"]

        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 200
        slides = resp.json()["extracted_data"]["slides"]

        assert slides[0]["notes"] == "These are speaker notes"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_get_extraction_after_extract(client, monkeypatch):
    """GET /api/curriculum/{id}/extraction returns the saved extraction."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        pptx_bytes = _make_simple_pptx()
        source = _upload_pptx(client, lesson_id, "saved.pptx", pptx_bytes)
        source_id = source["id"]

        # Extract first
        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 200

        # Now GET the saved extraction
        resp2 = client.get(f"/api/curriculum/{source_id}/extraction")
        assert resp2.status_code == 200
        body = resp2.json()
        assert body["id"] == source_id
        assert body["processing_status"] == "completed"
        assert body["extracted_data"] is not None
        assert body["extracted_data"]["source_type"] == "pptx"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_get_extraction_not_found(client):
    """GET extraction for a non-existent curriculum source returns 404."""
    random_id = "00000000-0000-0000-0000-000000000000"
    resp = client.get(f"/api/curriculum/{random_id}/extraction")
    assert resp.status_code == 404
    assert resp.json()["detail"] == "Curriculum source not found"


def test_extract_not_found(client):
    """POST extract for a non-existent curriculum source returns 404."""
    random_id = "00000000-0000-0000-0000-000000000000"
    resp = client.post(f"/api/curriculum/{random_id}/extract")
    assert resp.status_code == 404
    assert resp.json()["detail"] == "Curriculum source not found"


def test_extract_unsupported_file_type(client, monkeypatch):
    """Extracting a non-PPTX file returns 415."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        # Upload a PDF instead of PPTX
        resp_upload = client.post(
            f"/api/lessons/{lesson_id}/curriculum",
            files={"file": ("doc.pdf", io.BytesIO(b"%PDF-1.4 fake"), "application/pdf")},
        )
        assert resp_upload.status_code == 201
        source_id = resp_upload.json()["id"]

        # Attempt extraction
        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 415
        detail = resp.json()["detail"]
        assert ".pdf" in detail
        assert ".pptx" in detail
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_extract_status_transitions(client, monkeypatch):
    """Verify processing_status transitions: pending → processing → completed."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        pptx_bytes = _make_simple_pptx()
        source = _upload_pptx(client, lesson_id, "transitions.pptx", pptx_bytes)
        source_id = source["id"]

        # Should start as pending
        assert source["processing_status"] == "pending"

        # After extraction, should be completed
        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 200
        assert resp.json()["processing_status"] == "completed"

        # GET should also show completed
        resp2 = client.get(f"/api/curriculum/{source_id}/extraction")
        assert resp2.status_code == 200
        assert resp2.json()["processing_status"] == "completed"
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


def test_extract_preserves_slide_order(client, monkeypatch):
    """Verify that slides are returned in presentation order."""
    tmp = tempfile.mkdtemp()
    monkeypatch.setenv("STORAGE_ROOT", tmp)
    from app.core import config
    config.get_settings.cache_clear()

    try:
        lesson_id = _create_course_unit_lesson(client)
        pptx_bytes = _make_multi_slide_pptx()
        source = _upload_pptx(client, lesson_id, "order.pptx", pptx_bytes)
        source_id = source["id"]

        resp = client.post(f"/api/curriculum/{source_id}/extract")
        assert resp.status_code == 200
        slides = resp.json()["extracted_data"]["slides"]

        # Verify sequential slide numbers
        for i, slide in enumerate(slides, start=1):
            assert slide["slide_number"] == i
    finally:
        config.get_settings.cache_clear()
        shutil.rmtree(tmp, ignore_errors=True)


# ── Tests: PptxExtractor unit tests ─────────────────────────────


def test_pptx_extractor_direct():
    """Test PptxExtractor directly without the API layer."""
    from app.services.extractors import PptxExtractor

    pptx_bytes = _make_multi_slide_pptx()
    result = PptxExtractor().extract(pptx_bytes)

    assert result["source_type"] == "pptx"
    assert result["metadata"]["slide_count"] == 3
    assert len(result["slides"]) == 3
    assert result["slides"][0]["slide_number"] == 1
    assert result["slides"][2]["slide_number"] == 3


def test_pptx_extractor_empty_title():
    """Test that a slide with no title returns None for title."""
    from app.services.extractors import PptxExtractor

    prs = Presentation()
    # Use a blank layout (no title placeholder)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Add a text box
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2)).text_frame.text = "Text only"

    buf = io.BytesIO()
    prs.save(buf)

    result = PptxExtractor().extract(buf.getvalue())
    assert result["slides"][0]["title"] is None
    assert "Text only" in result["slides"][0]["texts"]
