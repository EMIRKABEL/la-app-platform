"""PPTX extraction service.

Uses ``python-pptx`` to read presentation structure and produce
a structured JSON representation containing slides, titles, text
blocks, tables, and speaker notes.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.extractors.base import BaseExtractor, ExtractionResult


class PptxExtractor(BaseExtractor):
    """Extract structured content from a ``.pptx`` file."""

    supported_extension = ".pptx"

    def extract(self, data: bytes) -> ExtractionResult:
        """Extract slides, titles, texts, tables, and notes from PPTX data.

        Args:
            data: Raw ``.pptx`` file bytes.

        Returns:
            An ``ExtractionResult`` with the following shape::

                {
                    "source_type": "pptx",
                    "metadata": {"slide_count": N, ...},
                    "slides": [
                        {
                            "slide_number": 1,
                            "title": "...",
                            "texts": ["...", ...],
                            "tables": [[["cell", ...], ...], ...],
                            "notes": "..." | null,
                        },
                        ...
                    ],
                }
        """
        presentation = Presentation(io.BytesIO(data))

        slides: list[dict[str, Any]] = []

        for index, slide in enumerate(presentation.slides, start=1):
            title = self._extract_title(slide)
            texts = self._extract_texts(slide)
            tables = self._extract_tables(slide)
            notes = self._extract_notes(slide)

            slides.append(
                {
                    "slide_number": index,
                    "title": title,
                    "texts": texts,
                    "tables": tables,
                    "notes": notes,
                }
            )

        metadata: dict[str, Any] = {
            "slide_count": len(slides),
        }

        return ExtractionResult(
            source_type="pptx",
            metadata=metadata,
            slides=slides,
        )

    # ── Helpers ──────────────────────────────────────────────────

    @staticmethod
    def _extract_title(slide: Any) -> str | None:
        """Return the slide title if it exists and is non-empty."""
        if slide.shapes.title is None:
            return None
        title_text = slide.shapes.title.text
        if title_text is None:
            return None
        title_text = title_text.strip()
        return title_text if title_text else None

    @staticmethod
    def _extract_texts(slide: Any) -> list[str]:
        """Extract all visible text block contents from a slide."""
        texts: list[str] = []
        for shape in slide.shapes:
            # Skip the title placeholder (already captured separately)
            if shape == slide.shapes.title:
                continue
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                para_text = "".join(run.text for run in paragraph.runs)
                if not para_text:
                    # Fall back to paragraph.text which may include
                    # text from runs not enumerated above
                    para_text = paragraph.text
                para_text = para_text.strip()
                if para_text:
                    texts.append(para_text)
        return texts

    @staticmethod
    def _extract_tables(slide: Any) -> list[list[list[str]]]:
        """Extract all table contents from a slide.

        Returns a list of tables, where each table is a list of rows,
        and each row is a list of cell string values.
        """
        tables: list[list[list[str]]] = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                continue
            table = shape.table
            rows: list[list[str]] = []
            for row in table.rows:
                cells: list[str] = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    cells.append(cell_text)
                rows.append(cells)
            tables.append(rows)
        return tables

    @staticmethod
    def _extract_notes(slide: Any) -> str | None:
        """Extract speaker notes from a slide if present."""
        if not slide.has_notes_slide:
            return None
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else None
        if notes_text is None:
            return None
        notes_text = notes_text.strip()
        return notes_text if notes_text else None
